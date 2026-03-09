"""
Carisma Wellness Group — Investor Deck: 7 New Slides
Generates a PowerPoint presentation matching the Carisma brand aesthetic.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Brand colours ───
BEIGE_BG = RGBColor(0xF5, 0xF0, 0xEB)       # Warm beige background
GOLD = RGBColor(0xB8, 0x9B, 0x72)            # Primary gold (headers, accents)
TAUPE = RGBColor(0x8C, 0x7B, 0x6B)           # Body text
DARK_TAUPE = RGBColor(0x5C, 0x4F, 0x45)      # Emphasis text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GOLD = RGBColor(0xD4, 0xC4, 0xA8)      # Subtle accents
ROSE_GOLD = RGBColor(0xC4, 0x8E, 0x7C)       # Secondary accent

# ─── Fonts ───
HEADER_FONT = "Didot"      # Elegant serif (fallback: Georgia)
BODY_FONT = "Avenir"       # Clean sans-serif (fallback: Calibri)

# ─── Slide dimensions (widescreen 16:9) ───
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background colour for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=11,
                 font_color=TAUPE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=BODY_FONT, spacing=1.15):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(spacing)
    return txBox, tf


def add_paragraph(tf, text, font_size=11, font_color=TAUPE, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name=BODY_FONT,
                  space_before=0, space_after=6):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def add_gold_line(slide, left, top, width):
    """Add a thin gold horizontal line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


def add_circle_metric(slide, center_x, top, number, label):
    """Add a circle with a large metric number and label below."""
    diameter = Inches(1.3)
    left = center_x - diameter // 2

    # Circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, diameter, diameter
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xEC, 0xE6, 0xDC)  # Light beige fill
    circle.line.color.rgb = LIGHT_GOLD
    circle.line.width = Pt(1)

    # Number inside circle
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(22)
    p.font.color.rgb = DARK_TAUPE
    p.font.name = HEADER_FONT
    p.font.bold = False
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(20)

    # Label below circle
    label_top = top + diameter + Inches(0.15)
    add_text_box(slide, left - Inches(0.3), label_top, diameter + Inches(0.6), Inches(0.5),
                 label, font_size=8, font_color=TAUPE, alignment=PP_ALIGN.CENTER,
                 font_name=BODY_FONT)


def add_gold_accent_bar(slide, left, top, width, height, text, font_size=10):
    """Add a gold accent bar with text."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xF0, 0xE6, 0xD5)
    bar.line.color.rgb = GOLD
    bar.line.width = Pt(0.5)

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = DARK_TAUPE
    p.font.name = BODY_FONT
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return bar


def add_metric_box(slide, left, top, width, height, number, label):
    """Add a small metric callout box."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xE6, 0xD5)
    box.line.color.rgb = GOLD
    box.line.width = Pt(0.5)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(20)
    p.font.color.rgb = GOLD
    p.font.name = HEADER_FONT
    p.font.bold = False
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(7)
    p2.font.color.rgb = TAUPE
    p2.font.name = BODY_FONT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(2)

    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return box


# ═══════════════════════════════════════════════════
# BUILD THE PRESENTATION
# ═══════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
blank_layout = prs.slide_layouts[6]  # Blank layout


# ───────────────────────────────────────────────────
# SLIDE 1: COMPANY AT A GLANCE
# ───────────────────────────────────────────────────
slide1 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide1, BEIGE_BG)

# Header
add_text_box(slide1, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.5),
             "C A R I S M A   W E L L N E S S   G R O U P",
             font_size=24, font_color=GOLD, alignment=PP_ALIGN.CENTER,
             font_name=HEADER_FONT)

# Gold line
add_gold_line(slide1, Inches(5.5), Inches(1.1), Inches(2.3))

# Subtitle
add_text_box(slide1, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.4),
             "A T   A   G L A N C E",
             font_size=14, font_color=TAUPE, alignment=PP_ALIGN.CENTER,
             font_name=HEADER_FONT)

# Description
add_text_box(slide1, Inches(2), Inches(2.0), Inches(9.3), Inches(0.8),
             "We design, build, and operate luxury wellness centres inside 4 and 5-star hotels across Malta, complemented by standalone medical aesthetics and weight management clinics.",
             font_size=11, font_color=TAUPE, alignment=PP_ALIGN.CENTER)

# Five metric circles
metrics = [
    ("€3M", "Revenue (2025)"),
    ("25%", "EBITDA Margin"),
    ("10", "Locations\nAcross Malta"),
    ("70", "Employees"),
    ("34", "Years\nOperating"),
]
start_x = Inches(1.8)
spacing_x = Inches(2.2)
for i, (num, label) in enumerate(metrics):
    add_circle_metric(slide1, start_x + spacing_x * i, Inches(3.2), num, label)

# Hotel logos text strip
add_text_box(slide1, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.4),
             "InterContinental  ·  Hyatt Regency  ·  AX Hotels  ·  Hugo's H Hotel  ·  Ramla Bay Resort  ·  Qawra Palace  ·  Kempinski",
             font_size=9, font_color=LIGHT_GOLD, alignment=PP_ALIGN.CENTER,
             font_name=BODY_FONT)

# Tagline
add_text_box(slide1, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.4),
             "Malta's Leading Vertically Integrated Wellness Operator",
             font_size=12, font_color=GOLD, alignment=PP_ALIGN.CENTER,
             font_name=HEADER_FONT, bold=False)


# ───────────────────────────────────────────────────
# SLIDE 2: BUSINESS MODEL
# ───────────────────────────────────────────────────
slide2 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide2, BEIGE_BG)

# Header
add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.5),
             "T H R E E   E N G I N E S   O F   G R O W T H",
             font_size=22, font_color=GOLD, alignment=PP_ALIGN.CENTER,
             font_name=HEADER_FONT)

add_gold_line(slide2, Inches(5.5), Inches(1.1), Inches(2.3))

# Left column header
add_text_box(slide2, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.35),
             "H O W   W E   O P E R A T E",
             font_size=11, font_color=GOLD, alignment=PP_ALIGN.LEFT,
             font_name=HEADER_FONT)

# Left column content
left_items = [
    ("Hotel Spa Management", "Revenue-share and management fee contracts with Malta's leading 4 and 5-star hotels. We design, build, staff, and operate — the hotel provides the location and guests. Currently 8 locations."),
    ("Medical Aesthetics", "Standalone Carisma Aesthetics clinic offering Botox, dermal fillers, and advanced skin treatments. Doctor-led, premium positioning. Launched 2024."),
    ("Medical Weight Loss", "Carisma Slimming clinic: CoolSculpting, EMSculpt NEO, VelaShape, and supervised weight loss programmes. Fastest-growing vertical. Launched 2025."),
]

y_pos = Inches(1.9)
for title, desc in left_items:
    _, tf = add_text_box(slide2, Inches(0.8), y_pos, Inches(5.5), Inches(1.0),
                         title, font_size=11, font_color=DARK_TAUPE, bold=True)
    add_paragraph(tf, desc, font_size=9, font_color=TAUPE, space_before=4)
    y_pos += Inches(1.1)

# Vertical divider
add_gold_line(slide2, Inches(6.6), Inches(1.4), Pt(1.5))
div = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(1.4), Pt(1.5), Inches(4.5))
div.fill.solid()
div.fill.fore_color.rgb = GOLD
div.line.fill.background()

# Right column header
add_text_box(slide2, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.35),
             "G R O W T H   P L A Y B O O K",
             font_size=11, font_color=GOLD, alignment=PP_ALIGN.LEFT,
             font_name=HEADER_FONT)

# Organic growth
_, tf_org = add_text_box(slide2, Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.3),
                         "Organic", font_size=11, font_color=DARK_TAUPE, bold=True)
organic_items = [
    "Win new hotel management contracts across Malta and Gozo",
    "Launch new service verticals within existing locations",
    "Cross-sell between brands: Spa → Aesthetics → Slimming",
    "Deepen wallet share with 20,000 active subscribers",
]
for item in organic_items:
    add_paragraph(tf_org, "·  " + item, font_size=9, font_color=TAUPE, space_before=3, space_after=3)

# Inorganic growth
_, tf_inorg = add_text_box(slide2, Inches(7.0), Inches(3.8), Inches(5.5), Inches(0.3),
                           "Inorganic (The Next Chapter)", font_size=11, font_color=DARK_TAUPE, bold=True)
inorganic_items = [
    "Acquire existing spa and clinic operations at attractive valuations",
    "Bolt on established client bases, locations, and talent",
    "Consolidate Malta's fragmented wellness market",
    "Apply our proven operational playbook to underperforming assets",
]
for item in inorganic_items:
    add_paragraph(tf_inorg, "·  " + item, font_size=9, font_color=TAUPE, space_before=3, space_after=3)

# Bottom accent bar
add_gold_accent_bar(slide2, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.5),
                    '"Grew from 3 locations to 10 in 3 years — all organic. Now ready to accelerate with acquisitions."',
                    font_size=10)


# ───────────────────────────────────────────────────
# SLIDE 3: GROWTH TIMELINE
# ───────────────────────────────────────────────────
slide3 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide3, BEIGE_BG)

add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.5),
             "O U R   G R O W T H   S T O R Y",
             font_size=22, font_color=GOLD, alignment=PP_ALIGN.CENTER,
             font_name=HEADER_FONT)

add_gold_line(slide3, Inches(5.5), Inches(1.0), Inches(2.3))

add_text_box(slide3, Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.3),
             "F R O M   F A M I L Y   B U S I N E S S   T O   P L A T F O R M",
             font_size=12, font_color=TAUPE, alignment=PP_ALIGN.CENTER,
             font_name=HEADER_FONT)

# Timeline line
timeline_y = Inches(2.5)
add_gold_line(slide3, Inches(0.6), timeline_y, Inches(12.1))

# Timeline entries
timeline = [
    ("1990", "Founded\nin Turkey", "50+ centres"),
    ("2000", "25 Centres", "New verticals"),
    ("2010", "Entered\nMalta", "First hammam"),
    ("2015", "3 Hotel\nSpas", "Malta footprint"),
    ("2022", "New CEO ★", "Mert Gulen"),
    ("2023", "€2.0M", "5 locations"),
    ("2024", "€2.5M", "8 locations"),
    ("2025", "€3.0M", "10 locations"),
    ("2026+", "Inorganic\nGrowth →", "Acquisitions"),
]

start_x = Inches(0.6)
entry_spacing = Inches(1.42)
dot_size = Inches(0.2)

for i, (year, main, sub) in enumerate(timeline):
    x = start_x + entry_spacing * i

    # Dot on timeline
    is_star = "★" in main
    dot_color = ROSE_GOLD if is_star else GOLD
    dot = slide3.shapes.add_shape(
        MSO_SHAPE.OVAL, x - dot_size // 2 + Inches(0.35), timeline_y - dot_size // 2, dot_size, dot_size
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()

    # Year above
    add_text_box(slide3, x - Inches(0.1), timeline_y - Inches(0.55), Inches(0.9), Inches(0.3),
                 year, font_size=9, font_color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=HEADER_FONT)

    # Main text below
    add_text_box(slide3, x - Inches(0.15), timeline_y + Inches(0.2), Inches(1.0), Inches(0.6),
                 main, font_size=8, font_color=DARK_TAUPE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Sub text
    add_text_box(slide3, x - Inches(0.15), timeline_y + Inches(0.7), Inches(1.0), Inches(0.4),
                 sub, font_size=7, font_color=TAUPE,
                 alignment=PP_ALIGN.CENTER)

# Sidebar box
sidebar_left = Inches(8.5)
sidebar = slide3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, sidebar_left, Inches(4.0), Inches(4.3), Inches(2.2)
)
sidebar.fill.solid()
sidebar.fill.fore_color.rgb = RGBColor(0xF0, 0xE6, 0xD5)
sidebar.line.color.rgb = GOLD
sidebar.line.width = Pt(0.5)

tf = sidebar.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Family business, new leadership."
p.font.size = Pt(10)
p.font.color.rgb = DARK_TAUPE
p.font.bold = True
p.font.name = BODY_FONT

add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "Founder Mustafa Gulen built 50+ wellness centres across Turkey and Malta over three decades. CEO Mert Gulen now leads the next chapter: transforming a family-operated chain into a scalable acquisition platform.",
              font_size=8, font_color=TAUPE, space_before=4)

# Revenue bars
bar_labels = [("2023", "€2.0M", 0.6), ("2024", "€2.5M", 0.75), ("2025", "€3.0M", 0.9)]
bar_y_start = Inches(4.3)
for i, (year, rev, width_factor) in enumerate(bar_labels):
    y = bar_y_start + Inches(0.55) * i
    # Year label
    add_text_box(slide3, Inches(0.8), y, Inches(0.6), Inches(0.35),
                 year, font_size=8, font_color=TAUPE, bold=True, alignment=PP_ALIGN.RIGHT)
    # Bar
    bar = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), y + Inches(0.05), Inches(3.5 * width_factor), Inches(0.25)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    # Revenue label
    add_text_box(slide3, Inches(1.5) + Inches(3.5 * width_factor) + Inches(0.1), y, Inches(0.8), Inches(0.35),
                 rev, font_size=8, font_color=DARK_TAUPE, bold=True)

add_text_box(slide3, Inches(1.5), Inches(6.0), Inches(3.5), Inches(0.3),
             "+50% in 3 years", font_size=8, font_color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER, font_name=HEADER_FONT)


# ───────────────────────────────────────────────────
# SLIDE 4: CASE STUDY — HUGO'S H HOTEL SPA
# ───────────────────────────────────────────────────
slide4 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide4, BEIGE_BG)

add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4),
             "C A S E   S T U D Y",
             font_size=14, font_color=TAUPE, alignment=PP_ALIGN.CENTER,
             font_name=HEADER_FONT)

add_text_box(slide4, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.5),
             "H U G O ' S   H   H O T E L   S P A",
             font_size=22, font_color=GOLD, alignment=PP_ALIGN.CENTER,
             font_name=HEADER_FONT)

add_gold_line(slide4, Inches(5.5), Inches(1.4), Inches(2.3))

# Hero IRR
add_text_box(slide4, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.0),
             "88% IRR",
             font_size=60, font_color=GOLD, alignment=PP_ALIGN.CENTER,
             font_name=HEADER_FONT)

add_text_box(slide4, Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.4),
             "Built and launched a luxury spa during a global pandemic.",
             font_size=13, font_color=TAUPE, alignment=PP_ALIGN.CENTER,
             font_name=BODY_FONT, bold=False)

# Three columns
col_width = Inches(3.6)
col_gap = Inches(0.3)
col_start = Inches(0.8)
col_top = Inches(3.5)

columns = [
    ("T H E   C H A L L E N G E",
     "Won the Hugo's H Hotel spa management contract at the height of COVID-19.\n\nEmpty hotel space. No spa infrastructure. Supply chain disruptions. Construction crew limitations. Import restrictions on specialist equipment.\n\nTimeline pressure: hotel opening date was fixed regardless of COVID delays."),
    ("O U R   E X E C U T I O N",
     "Full turnkey delivery in-house: concept design, procurement through our established supplier network, construction management, recruitment, and staff training.\n\nLeveraged 30+ years of vertical integration and relationships with global suppliers to navigate pandemic bottlenecks.\n\nDelivered on time, on budget."),
    ("T H E   R E S U L T",
     "Fully operational luxury spa delivered in 5 months despite pandemic conditions.\n\n88% internal rate of return.\n\nNow one of our highest-performing locations, contributing significantly to group revenue and demonstrating the Carisma operational playbook under the most challenging conditions imaginable."),
]

for i, (header, body) in enumerate(columns):
    x = col_start + (col_width + col_gap) * i

    # Column header
    add_text_box(slide4, x, col_top, col_width, Inches(0.3),
                 header, font_size=9, font_color=GOLD, bold=False,
                 font_name=HEADER_FONT)
    add_gold_line(slide4, x, col_top + Inches(0.3), col_width)

    # Column body
    add_text_box(slide4, x, col_top + Inches(0.45), col_width, Inches(2.3),
                 body, font_size=8.5, font_color=TAUPE)

# Bottom callout
add_gold_accent_bar(slide4, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
                    '"If we can deliver 88% returns building a spa during a pandemic, imagine the returns when we acquire underperforming operations and apply the same playbook."',
                    font_size=9)


# ───────────────────────────────────────────────────
# SLIDE 5: CAPITAL STRUCTURE — THE ASK
# ───────────────────────────────────────────────────
slide5 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide5, BEIGE_BG)

add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.5),
             "C A P I T A L   S T R U C T U R E",
             font_size=22, font_color=GOLD, alignment=PP_ALIGN.CENTER,
             font_name=HEADER_FONT)

add_gold_line(slide5, Inches(5.5), Inches(1.0), Inches(2.3))

add_text_box(slide5, Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.3),
             "A C C E L E R A T I N G   I N O R G A N I C   G R O W T H",
             font_size=12, font_color=TAUPE, alignment=PP_ALIGN.CENTER,
             font_name=HEADER_FONT)

# Left panel — Today
panel_top = Inches(1.8)
add_text_box(slide5, Inches(0.8), panel_top, Inches(3.5), Inches(0.3),
             "T O D A Y", font_size=11, font_color=GOLD, bold=False,
             font_name=HEADER_FONT)
add_gold_line(slide5, Inches(0.8), panel_top + Inches(0.35), Inches(3.5))

today_items = [
    "€3.0M Revenue",
    "€750K EBITDA (25% margin)",
    "10 Locations",
    "70 Employees",
    "Zero Debt",
    "",
    "Full asset base:",
    "· Equipment across 10 locations",
    "· Long-term hotel contracts",
    "· 20,000-subscriber client network",
    "· 3 established brands",
    "· 34-year operational playbook",
    "· Distribution rights",
]

_, tf_today = add_text_box(slide5, Inches(0.8), panel_top + Inches(0.5), Inches(3.5), Inches(4.0),
                           today_items[0], font_size=9, font_color=DARK_TAUPE, bold=True)
for item in today_items[1:]:
    is_bold = bool(item and not item.startswith("·") and item != "Full asset base:")
    add_paragraph(tf_today, item, font_size=9,
                  font_color=DARK_TAUPE if is_bold else TAUPE,
                  bold=is_bold, space_before=2, space_after=2)

# Center panel — NewCo Structure
center_x = Inches(4.8)
add_text_box(slide5, center_x, panel_top, Inches(3.7), Inches(0.3),
             "P R O P O S E D   S T R U C T U R E", font_size=11, font_color=GOLD,
             font_name=HEADER_FONT)
add_gold_line(slide5, center_x, panel_top + Inches(0.35), Inches(3.7))

# NewCo box
newco_box = slide5.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, center_x, panel_top + Inches(0.5), Inches(3.7), Inches(3.8)
)
newco_box.fill.solid()
newco_box.fill.fore_color.rgb = RGBColor(0xF0, 0xE6, 0xD5)
newco_box.line.color.rgb = GOLD
newco_box.line.width = Pt(1)

tf = newco_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "CARISMA NEWCO"
p.font.size = Pt(13)
p.font.color.rgb = GOLD
p.font.name = HEADER_FONT
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(12)

sections = [
    ("Carisma Existing Assets", "(Founder contribution)", "Valuation: [TBD]"),
    ("", "", ""),
    ("Growth Capital: [TBD]", "", ""),
    ("Debt [TBD]  |  Equity [TBD]", "", ""),
    ("", "", ""),
    ("Target Leverage: ≤ 2x EBITDA", "(Conservative, asset-backed)", ""),
]
for main, sub1, sub2 in sections:
    if main:
        add_paragraph(tf, main, font_size=10, font_color=DARK_TAUPE, bold=True,
                      alignment=PP_ALIGN.CENTER, space_before=6, space_after=2)
    if sub1:
        add_paragraph(tf, sub1, font_size=8, font_color=TAUPE,
                      alignment=PP_ALIGN.CENTER, space_before=0, space_after=2)
    if sub2:
        add_paragraph(tf, sub2, font_size=9, font_color=TAUPE,
                      alignment=PP_ALIGN.CENTER, space_before=0, space_after=4)

# Right panel — Use of Funds
right_x = Inches(9.0)
add_text_box(slide5, right_x, panel_top, Inches(3.7), Inches(0.3),
             "U S E   O F   F U N D S", font_size=11, font_color=GOLD,
             font_name=HEADER_FONT)
add_gold_line(slide5, right_x, panel_top + Inches(0.35), Inches(3.7))

funds_items = [
    ("2–3 acquisitions", "over the next 24 months"),
    ("Working capital", "for post-acquisition integration and operational improvement"),
    ("Clinic buildout capex", "within acquired locations (aesthetics + slimming bolt-ons)"),
    ("Market consolidation", "build Malta's dominant wellness platform"),
]

y = panel_top + Inches(0.5)
for title, desc in funds_items:
    _, tf_f = add_text_box(slide5, right_x, y, Inches(3.7), Inches(0.8),
                           title, font_size=10, font_color=DARK_TAUPE, bold=True)
    add_paragraph(tf_f, desc, font_size=8, font_color=TAUPE, space_before=2)
    y += Inches(0.85)

# Footer
add_text_box(slide5, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
             "Structure and quantum being finalised with M&A advisor. Conservative leverage with clear path to deleveraging through EBITDA growth. Founder retains majority ownership and operational control.",
             font_size=7, font_color=LIGHT_GOLD, alignment=PP_ALIGN.CENTER)


# ───────────────────────────────────────────────────
# SLIDE 6: ACQUISITION PIPELINE
# ───────────────────────────────────────────────────
slide6 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide6, BEIGE_BG)

add_text_box(slide6, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.5),
             "A C Q U I S I T I O N   P I P E L I N E",
             font_size=22, font_color=GOLD, alignment=PP_ALIGN.CENTER,
             font_name=HEADER_FONT)

add_gold_line(slide6, Inches(5.5), Inches(1.0), Inches(2.3))

add_text_box(slide6, Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.3),
             "T W O   I D E N T I F I E D   T A R G E T S",
             font_size=12, font_color=TAUPE, alignment=PP_ALIGN.CENTER,
             font_name=HEADER_FONT)

# Target A card
card_top = Inches(1.7)
card_height = Inches(4.2)
card_a = slide6.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), card_top, Inches(5.6), card_height
)
card_a.fill.solid()
card_a.fill.fore_color.rgb = WHITE
card_a.line.color.rgb = GOLD
card_a.line.width = Pt(0.75)

# Target A content
add_text_box(slide6, Inches(1.1), card_top + Inches(0.15), Inches(5.0), Inches(0.3),
             "INA SPA & WELLNESS", font_size=13, font_color=GOLD,
             font_name=HEADER_FONT)
add_text_box(slide6, Inches(1.1), card_top + Inches(0.5), Inches(5.0), Inches(0.2),
             "Numi Hotel, Malta", font_size=9, font_color=TAUPE, bold=False)

metrics_a = [
    ("Revenue", "~€360K/year"),
    ("EBITDA", "~€72K/year"),
    ("Target Price", "€300–400K"),
    ("Expected IRR", "36–40%"),
    ("Locations", "1 hotel spa"),
    ("Staff", "14 employees"),
]
y = card_top + Inches(0.85)
for label, value in metrics_a:
    add_text_box(slide6, Inches(1.1), y, Inches(2.0), Inches(0.2),
                 label, font_size=8, font_color=TAUPE)
    add_text_box(slide6, Inches(3.1), y, Inches(3.0), Inches(0.2),
                 value, font_size=8, font_color=DARK_TAUPE, bold=True)
    y += Inches(0.22)

add_text_box(slide6, Inches(1.1), y + Inches(0.1), Inches(5.0), Inches(0.2),
             "Rationale", font_size=9, font_color=GOLD, bold=False, font_name=HEADER_FONT)
add_text_box(slide6, Inches(1.1), y + Inches(0.35), Inches(5.0), Inches(0.8),
             "Motivated seller with breakeven operations due to overstaffing (€88–110K excess labour). Carisma's playbook immediately cuts costs and bolts on a medical aesthetics clinic — proven model.",
             font_size=8, font_color=TAUPE)

add_text_box(slide6, Inches(1.1), y + Inches(1.1), Inches(5.0), Inches(0.2),
             "Seller: Distressed. Invested €1M, runs at breakeven. Ready to exit.",
             font_size=7.5, font_color=ROSE_GOLD, bold=False)

# Target B card
card_b = slide6.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(6.9), card_top, Inches(5.6), card_height
)
card_b.fill.solid()
card_b.fill.fore_color.rgb = WHITE
card_b.line.color.rgb = GOLD
card_b.line.width = Pt(0.75)

add_text_box(slide6, Inches(7.2), card_top + Inches(0.15), Inches(5.0), Inches(0.3),
             "MARION MIZZI WELLBEING", font_size=13, font_color=GOLD,
             font_name=HEADER_FONT)
add_text_box(slide6, Inches(7.2), card_top + Inches(0.5), Inches(5.0), Inches(0.2),
             "Malta's Market Leader in Slimming & Weight Loss", font_size=9, font_color=TAUPE)

metrics_b = [
    ("Revenue", "~€2.0M/year"),
    ("EBITDA", "~€500K/year"),
    ("Target Price", "€1.5–2.0M"),
    ("Expected IRR", "18–24%"),
    ("Locations", "3 (Fgura, Sliema, Mellieha)"),
    ("Client Base", "4,900+ active clients"),
    ("Market Share", "~35% (weight management)"),
    ("Brand History", "~50 years established"),
]
y = card_top + Inches(0.85)
for label, value in metrics_b:
    add_text_box(slide6, Inches(7.2), y, Inches(2.0), Inches(0.2),
                 label, font_size=8, font_color=TAUPE)
    add_text_box(slide6, Inches(9.2), y, Inches(3.0), Inches(0.2),
                 value, font_size=8, font_color=DARK_TAUPE, bold=True)
    y += Inches(0.22)

add_text_box(slide6, Inches(7.2), y + Inches(0.1), Inches(5.0), Inches(0.2),
             "Rationale", font_size=9, font_color=GOLD, bold=False, font_name=HEADER_FONT)
add_text_box(slide6, Inches(7.2), y + Inches(0.35), Inches(5.0), Inches(0.8),
             "Malta's dominant weight management brand. ~35% market share, 50 years of brand equity. Creates near-monopoly (85–95% combined share). Massive cross-sell into Carisma Aesthetics and Spa.",
             font_size=8, font_color=TAUPE)

add_text_box(slide6, Inches(7.2), y + Inches(1.1), Inches(5.0), Inches(0.2),
             "Seller: Succession-driven. Founder retired, daughter disengaged. Open to exit.",
             font_size=7.5, font_color=ROSE_GOLD, bold=False)

# Combined impact bar
add_gold_accent_bar(slide6, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.55),
                    "Post-acquisition combined entity:  €5.3M+ revenue  |  13+ locations  |  Dominant market position across spa, aesthetics, and weight management in Malta",
                    font_size=10)


# ───────────────────────────────────────────────────
# SLIDE 7: WHY NOW — INVESTMENT THESIS
# ───────────────────────────────────────────────────
slide7 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide7, BEIGE_BG)

add_text_box(slide7, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.5),
             "W H Y   N O W",
             font_size=28, font_color=GOLD, alignment=PP_ALIGN.CENTER,
             font_name=HEADER_FONT)

add_gold_line(slide7, Inches(5.5), Inches(1.1), Inches(2.3))

# Six thesis bullets
theses = [
    ("1.", "Proven Operator, Proven Growth",
     "50% revenue growth in 3 years. 25% EBITDA margins. 88% IRR on latest buildout. All organic. Management has demonstrated it can build and operate profitably before asking for acquisition capital."),
    ("2.", "Fragmented Market Ripe for Consolidation",
     "Malta's wellness market has no dominant platform player. Dozens of small, single-location operators with no scale advantages. Classic roll-up opportunity in a growing market."),
    ("3.", "Motivated Sellers at Attractive Valuations",
     "Both pipeline targets are distress or succession-driven, enabling acquisitions at 3–4x EBITDA versus 6–8x for healthy businesses. Built-in margin of safety."),
    ("4.", "Vertical Integration = Instant Value Creation",
     "We design, build, procure, and operate. Every acquisition immediately benefits from our supply chain, training systems, and 34-year operational playbook. Day-one cost reduction. Day-one cross-sell across three brands."),
    ("5.", "Multiple Expansion Potential",
     "Acquiring single-location operators at 3–4x and integrating into a multi-location platform commanding higher multiples. Each bolt-on is immediately accretive to group value."),
    ("6.", "Malta Market Tailwinds",
     "530,000 population growing 2–3% annually. 2.8M+ tourists per year, growing 3–4%. Highest obesity rate in Europe drives slimming demand. Wellness tourism is a global megatrend."),
]

y = Inches(1.5)
for num, title, desc in theses:
    # Number
    add_text_box(slide7, Inches(0.8), y, Inches(0.35), Inches(0.3),
                 num, font_size=11, font_color=GOLD, bold=True, font_name=HEADER_FONT)
    # Title
    add_text_box(slide7, Inches(1.15), y, Inches(8.0), Inches(0.25),
                 title, font_size=10, font_color=DARK_TAUPE, bold=True)
    # Description
    add_text_box(slide7, Inches(1.15), y + Inches(0.25), Inches(8.0), Inches(0.5),
                 desc, font_size=8, font_color=TAUPE)
    y += Inches(0.82)

# Right margin callout boxes
callouts = [
    ("8", "Industry\nAwards"),
    ("99%", "Customer\nSatisfaction"),
    ("34", "Years\nOperating"),
    ("20,000", "Active\nSubscribers"),
]

box_top = Inches(1.5)
for i, (number, label) in enumerate(callouts):
    add_metric_box(slide7, Inches(10.2), box_top + Inches(1.2) * i,
                   Inches(2.3), Inches(0.95), number, label)

# Closing statement
add_text_box(slide7, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.5),
             '"We\'ve spent 34 years building the operational playbook. Now we\'re ready to deploy it through acquisitions — and the market timing has never been better."',
             font_size=10, font_color=GOLD, alignment=PP_ALIGN.CENTER,
             font_name=HEADER_FONT, bold=False)


# ═══════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════
output_dir = "/Users/mertgulen/Library/CloudStorage/GoogleDrive-mertgulen98@gmail.com/My Drive/Carisma Wellness Group/Carisma AI /Carisma AI/.tmp"
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "Carisma-Investor-Deck-7-New-Slides.pptx")
prs.save(output_path)
print(f"Saved to: {output_path}")
